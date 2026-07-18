import json
import os


def read_presentation(args_json):
    args = json.loads(args_json)
    try:
        from pptx import Presentation
        prs = Presentation(args["input_path"])
        results = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.extend([p.text for p in shape.text_frame.paragraphs if p.text.strip()])
            results.append({"slide": i + 1, "texts": texts})
        return json.dumps(results, indent=2)
    except ImportError:
        return "ERROR: python-pptx not available"


def create_presentation(args_json):
    args = json.loads(args_json)
    slides_data = json.loads(args.get("slides", "[]"))
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation()
        for slide_info in slides_data:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            body = slide.placeholders[1]
            title.text = slide_info.get("title", "")
            tf = body.text_frame
            tf.text = ""
            for point in slide_info.get("content", []):
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
        prs.save(args["output_path"])
        return f"Created presentation with {len(slides_data)} slides → {args['output_path']}"
    except ImportError:
        return "ERROR: python-pptx not available"


def presentation_to_pdf(args_json):
    args = json.loads(args_json)
    try:
        from pptx import Presentation
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import landscape
        prs = Presentation(args["input_path"])
        c = canvas.Canvas(args["output_path"], pagesize=landscape(prs.slide_width, prs.slide_height))
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            c.drawString(shape.left, shape.top, para.text[:100])
            c.showPage()
        c.save()
        return f"Converted PPTX → PDF: {args['output_path']}"
    except Exception as e:
        return f"ERROR: {e}"
